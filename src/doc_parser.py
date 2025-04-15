import pdfplumber
import re
import os
import logging
from typing import Optional, List
from docx import Document
from pptx import Presentation

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


import logging
from typing import Optional

try:
    import pdfplumber
except ImportError:
    pdfplumber = None
    logging.warning("pdfplumber not installed. PDF functionality will be limited.")

try:
    from docx import Document
except ImportError:
    Document = None
    logging.warning("python-docx not installed. DOCX functionality will be limited.")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    logging.warning("python-pptx not installed. PPTX functionality will be limited.")


def extract_text_from_file(file_path: str) -> Optional[str]:
    """Extracts text content from PDF, DOCX, PPTX, or TXT files."""
    try:
        file_lower = file_path.lower()
        if file_lower.endswith(".pdf"):
            if pdfplumber:
                with pdfplumber.open(file_path) as pdf:
                    text_content = []
                    for page_number, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text:
                            text_content.append(text)
                        else:
                            logging.warning(f"No text found in {file_path}, page {page_number + 1}")
                    return "\n".join(text_content)
            else:
                logging.error("pdfplumber is not installed. Cannot process PDF files.")
                return None
        elif file_lower.endswith(".docx"):
            if Document:
                doc = Document(file_path)
                text_content = [paragraph.text for paragraph in doc.paragraphs]
                return "\n".join(text_content)
            else:
                logging.error("python-docx is not installed. Cannot process DOCX files.")
                return None
        elif file_lower.endswith(".pptx"):
            if Presentation:
                prs = Presentation(file_path)
                text_content = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content.append(shape.text)
                return "\n".join(text_content)
            else:
                logging.error("python-pptx is not installed. Cannot process PPTX files.")
                return None
        elif file_lower.endswith(".txt"):
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            logging.error(f"Unsupported file type: {file_path}")
            return None
    except FileNotFoundError:
        logging.error(f"File not found: {file_path}")
        return None
    except Exception as e:
        logging.exception(f"Error reading {file_path}: {e}")
        return None


def clean_text(text: str) -> str:
    """Cleans the text by removing extra whitespace."""
    cleaned_text = re.sub(r'\s+', ' ', text)
    return cleaned_text.strip()


def process_folder(folder_path: str, output_directory: str) -> None:
    """Processes all compatible files within a single folder and writes to a separate file."""
    folder_name = os.path.basename(folder_path)
    output_filename = f"{folder_name}_processed.txt"
    output_path = os.path.join(output_directory, output_filename)
    all_text_in_folder = []

    logging.info(f"Processing folder: {folder_path}")
    for filename in os.listdir(folder_path):
        if filename.lower().endswith((".pdf", ".docx", ".pptx",".txt")):
            file_path = os.path.join(folder_path, filename)
            logging.info(f"Processing file: {filename}")
            raw_text = extract_text_from_file(file_path)
            if raw_text:
                cleaned_text = clean_text(raw_text)
                text_with_metadata = f"File: {filename}\n\n{cleaned_text}"
                all_text_in_folder.append(text_with_metadata[0:127999])
            else:
                logging.warning(f"Failed to extract text from {file_path}")

    if all_text_in_folder:
        try:
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("\n\n".join(all_text_in_folder))
            logging.info(f"Successfully wrote text from '{folder_name}' to {output_path}")
        except Exception as e:
            logging.exception(f"Error writing to {output_path}: {e}")
    else:
        logging.warning(f"No text was produced in '{folder_name}', so no file was written.")


def process_folders(root_folder: str, output_root: str) -> None:
    """Iterates through subfolders in the root folder and processes each one."""
    os.makedirs(output_root, exist_ok=True)  # Ensure the output directory exists
    for item in os.listdir(root_folder):
        item_path = os.path.join(root_folder, item)
        if os.path.isdir(item_path):
            process_folder(item_path, output_root)



    
# Example Usage:
if __name__ == "__main__":
    # Step 1: Extract text from documents in subfolders
    # root_directory = r"C:\SAFELITE\SAFELITE\data"  # Raw string for Windows paths
    # processed_data_dir = r"C:\SAFELITE\SAFELITE\processed_data_by_folder"

    # Get the directory of the current script (the src folder)
    src_dir = os.path.dirname(os.path.abspath(__file__))

    # Construct the path to the root directory (one level up from src)
    root_directory = os.path.dirname(src_dir)

    # Construct the path to the data directory
    data_dir = os.path.join(root_directory, "data")

    # Construct the path to the processed_data_by_folder directory
    processed_data_dir = os.path.join(root_directory, "processed_data_by_folder")

    # Process all subfolders and extract text
    process_folders(root_directory, processed_data_dir)

